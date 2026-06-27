"""Text extraction and natural-boundary chunking."""

from __future__ import annotations

import ast
import io
import json
import re
import shutil
import sqlite3
import traceback
from collections import Counter
from collections.abc import Iterable, Mapping
from contextlib import closing
from dataclasses import dataclass
from datetime import datetime, timezone
from importlib import metadata
from pathlib import Path
from typing import Any

from uni_rag_agent.config import Config
from uni_rag_agent.storage import connect_sqlite, ensure_data_dirs, initialize_schema

TEXT_EXTRACTION_CATEGORIES = {
    "document",
    "slides",
    "notebook",
    "code",
    "transcript",
}

SUPPORTED_TEXT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".pptx",
    ".ipynb",
    ".py",
    ".r",
    ".cpp",
    ".h",
    ".m",
    ".vtt",
}

LEGACY_EXTENSIONS = {".doc", ".ppt"}
LEGACY_FORMAT_REASON = "legacy format not supported yet"
SCANNED_PDF_OCR_REASON = "scanned PDF, OCR not available"
NO_TEXT_REASON = "no extractable text found"

DEFAULT_MAX_CHUNK_TOKENS = 1000
NOTEBOOK_OUTPUT_CHAR_LIMIT = 500
ERROR_CHAR_LIMIT = 4000
PDF_SCANNED_TEXT_CHAR_THRESHOLD = 20

TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "cp1252", "latin-1")
VTT_TIMESTAMP_RE = re.compile(
    r"^(?P<start>\d{2}:\d{2}(?::\d{2})?(?:\.\d{3})?)\s+-->\s+"
    r"(?P<end>\d{2}:\d{2}(?::\d{2})?(?:\.\d{3})?)"
)
R_FUNCTION_RE = re.compile(r"^\s*([A-Za-z.][\w.]*)\s*(?:<-|=)\s*function\s*\(", re.M)
CPP_FUNCTION_RE = re.compile(
    r"^\s*(?:[\w:<>,~*&]+\s+)+([A-Za-z_]\w*)\s*\([^;{}]*\)\s*(?:const\s*)?\{",
    re.M,
)
MATLAB_FUNCTION_RE = re.compile(
    r"^\s*function\s+(?:\[[^\]]+\]\s*=\s*|[A-Za-z_]\w*\s*=\s*)?"
    r"(?P<name>[A-Za-z_]\w*)",
    re.M,
)


class ExtractionError(RuntimeError):
    """Raised when an extraction run cannot be completed."""


class ExtractionFailure(RuntimeError):
    """Raised for a single file that cannot be extracted."""

    def __init__(
        self,
        message: str,
        *,
        extractor_name: str,
        extractor_version: str | None = None,
        metadata_json: str | None = None,
    ) -> None:
        super().__init__(message)
        self.extractor_name = extractor_name
        self.extractor_version = extractor_version
        self.metadata_json = metadata_json


@dataclass(frozen=True)
class PendingFileRecord:
    id: int
    path: Path
    relative_path: str
    filename: str
    extension: str
    category: str
    content_hash: str | None


@dataclass(frozen=True)
class RawChunk:
    source_type: str
    title: str | None
    text: str
    location_type: str
    location_value: str
    metadata: Mapping[str, object]


@dataclass(frozen=True)
class ChunkRecord:
    file_id: int
    chunk_uid: str
    source_type: str
    chunk_index: int
    title: str | None
    text: str
    token_count: int
    location_type: str
    location_value: str
    metadata_json: str


@dataclass(frozen=True)
class ExtractedDocument:
    file_id: int
    extractor_name: str
    extractor_version: str | None
    status: str
    text_length: int
    chunk_count: int
    metadata_json: str
    error: str | None
    chunks: tuple[ChunkRecord, ...]


@dataclass(frozen=True)
class ExtractionFailureSummary:
    file_id: int
    path: str
    error: str


@dataclass(frozen=True)
class ExtractionRunResult:
    run_id: int
    started_at: str
    finished_at: str
    status: str
    category: str | None
    files_seen: int
    files_indexed: int
    files_failed: int
    chunks_created: int
    by_source_type: Mapping[str, int]
    failures: tuple[ExtractionFailureSummary, ...]
    diagnostics: tuple[str, ...]

    def as_safe_dict(self) -> dict[str, object]:
        return {
            "run_id": self.run_id,
            "started_at": self.started_at,
            "finished_at": self.finished_at,
            "status": self.status,
            "category": self.category,
            "files_seen": self.files_seen,
            "files_indexed": self.files_indexed,
            "files_failed": self.files_failed,
            "chunks_created": self.chunks_created,
            "by_source_type": dict(self.by_source_type),
            "failures": [failure.__dict__ for failure in self.failures],
            "diagnostics": list(self.diagnostics),
        }


@dataclass(frozen=True)
class ExtractionStatus:
    latest_extraction_run_id: int | None
    latest_extraction_started_at: str | None
    pending_text_files: int
    indexed_text_files: int
    failed_text_files: int
    extracted_documents: int
    chunks_total: int
    pending_by_category: Mapping[str, int]
    chunks_by_source_type: Mapping[str, int]
    recent_failures: tuple[ExtractionFailureSummary, ...]

    def as_safe_dict(self) -> dict[str, object]:
        return {
            "latest_extraction_run_id": self.latest_extraction_run_id,
            "latest_extraction_started_at": self.latest_extraction_started_at,
            "pending_text_files": self.pending_text_files,
            "indexed_text_files": self.indexed_text_files,
            "failed_text_files": self.failed_text_files,
            "extracted_documents": self.extracted_documents,
            "chunks_total": self.chunks_total,
            "pending_by_category": dict(self.pending_by_category),
            "chunks_by_source_type": dict(self.chunks_by_source_type),
            "recent_failures": [failure.__dict__ for failure in self.recent_failures],
        }


def extract_pending_files(
    config: Config,
    category: str | None = None,
) -> ExtractionRunResult:
    """Extract pending text-like files and persist chunks per file."""
    _validate_category(category)
    ensure_data_dirs(config)

    with closing(connect_sqlite(config)) as connection:
        initialize_schema(connection)
        started_at = _utc_now()
        run_id = _start_extraction_run(connection, config, category, started_at)
        connection.commit()

        pending_files = _load_pending_files(connection, category)
        files_indexed = 0
        files_failed = 0
        chunks_created = 0
        by_source_type: Counter[str] = Counter()
        failures: list[ExtractionFailureSummary] = []
        diagnostics: list[str] = []

        try:
            for file_record in pending_files:
                try:
                    extracted = extract_file(file_record, config)
                except ExtractionFailure as exc:
                    files_failed += 1
                    failure = _failure_from_exception(file_record, exc)
                    failures.append(failure)
                    _record_failed_extraction(
                        connection,
                        run_id=run_id,
                        file_record=file_record,
                        extractor_name=exc.extractor_name,
                        extractor_version=exc.extractor_version,
                        metadata_json=exc.metadata_json,
                        error=failure.error,
                    )
                    continue
                except Exception as exc:
                    files_failed += 1
                    error = _format_exception(exc)
                    failures.append(
                        ExtractionFailureSummary(
                            file_id=file_record.id,
                            path=str(file_record.path),
                            error=error,
                        )
                    )
                    _record_failed_extraction(
                        connection,
                        run_id=run_id,
                        file_record=file_record,
                        extractor_name=_extractor_name_for_extension(
                            file_record.extension
                        ),
                        extractor_version=None,
                        metadata_json=_json_dumps({"extension": file_record.extension}),
                        error=error,
                    )
                    continue

                files_indexed += 1
                chunks_created += extracted.chunk_count
                by_source_type.update(chunk.source_type for chunk in extracted.chunks)
                with connection:
                    _persist_successful_extraction(
                        connection,
                        run_id=run_id,
                        extracted=extracted,
                        extracted_at=_utc_now(),
                    )

            finished_at = _utc_now()
            status = "completed"
            with connection:
                _finish_extraction_run(
                    connection,
                    run_id=run_id,
                    finished_at=finished_at,
                    status=status,
                    files_seen=len(pending_files),
                    files_indexed=files_indexed,
                    files_failed=files_failed,
                    error=None,
                )
        except Exception as exc:
            with connection:
                _finish_extraction_run(
                    connection,
                    run_id=run_id,
                    finished_at=_utc_now(),
                    status="failed",
                    files_seen=len(pending_files),
                    files_indexed=files_indexed,
                    files_failed=files_failed,
                    error=str(exc),
                )
            raise

    return ExtractionRunResult(
        run_id=run_id,
        started_at=started_at,
        finished_at=finished_at,
        status=status,
        category=category,
        files_seen=len(pending_files),
        files_indexed=files_indexed,
        files_failed=files_failed,
        chunks_created=chunks_created,
        by_source_type=dict(sorted(by_source_type.items())),
        failures=tuple(failures),
        diagnostics=tuple(diagnostics),
    )


def load_extraction_status(config: Config) -> ExtractionStatus:
    """Read extraction/chunk coverage without traversing Courses."""
    if not config.sqlite_path.is_file():
        raise ExtractionError(f"SQLite database does not exist: {config.sqlite_path}")

    with closing(connect_sqlite(config)) as connection:
        latest_run_id, latest_started_at = _latest_extraction_run(connection)
        pending_by_category = _count_pending_by_category(connection)
        chunks_by_source_type = _count_chunks_by_source_type(connection)
        pending_text_files = sum(pending_by_category.values())
        indexed_text_files = _count_files_by_status(connection, "indexed")
        failed_text_files = _count_files_by_status(connection, "failed")
        extracted_documents = _count_table_rows(connection, "extracted_documents")
        chunks_total = _count_table_rows(connection, "chunks")
        recent_failures = _load_recent_failures(connection)

    return ExtractionStatus(
        latest_extraction_run_id=latest_run_id,
        latest_extraction_started_at=latest_started_at,
        pending_text_files=pending_text_files,
        indexed_text_files=indexed_text_files,
        failed_text_files=failed_text_files,
        extracted_documents=extracted_documents,
        chunks_total=chunks_total,
        pending_by_category=pending_by_category,
        chunks_by_source_type=chunks_by_source_type,
        recent_failures=recent_failures,
    )


def extract_file(file_record: PendingFileRecord, config: Config) -> ExtractedDocument:
    """Extract one pending file into final chunk records."""
    extension = file_record.extension
    extractor_name = _extractor_name_for_extension(extension)
    extractor_version = _extractor_version_for_extension(extension)

    if extension in LEGACY_EXTENSIONS:
        raise ExtractionFailure(
            LEGACY_FORMAT_REASON,
            extractor_name=extractor_name,
            extractor_version=extractor_version,
            metadata_json=_json_dumps({"extension": extension}),
        )

    raw_chunks = _extract_raw_chunks(file_record, config)
    chunks = _finalize_chunks(
        file_record=file_record,
        raw_chunks=raw_chunks,
        max_tokens=DEFAULT_MAX_CHUNK_TOKENS,
    )

    if not chunks:
        raise ExtractionFailure(
            NO_TEXT_REASON,
            extractor_name=extractor_name,
            extractor_version=extractor_version,
            metadata_json=_json_dumps({"extension": extension}),
        )

    metadata_json = _json_dumps(
        {
            "extension": extension,
            "relative_path": file_record.relative_path,
            "content_hash": file_record.content_hash,
            "max_chunk_tokens": DEFAULT_MAX_CHUNK_TOKENS,
        }
    )
    return ExtractedDocument(
        file_id=file_record.id,
        extractor_name=extractor_name,
        extractor_version=extractor_version,
        status="indexed",
        text_length=sum(len(chunk.text) for chunk in chunks),
        chunk_count=len(chunks),
        metadata_json=metadata_json,
        error=None,
        chunks=tuple(chunks),
    )


def _extract_raw_chunks(
    file_record: PendingFileRecord,
    config: Config,
) -> tuple[RawChunk, ...]:
    extension = file_record.extension
    path = file_record.path
    if extension == ".pdf":
        return _extract_pdf(path, config)
    if extension == ".pptx":
        return _extract_pptx(path)
    if extension == ".docx":
        return _extract_docx(path)
    if extension == ".txt":
        return _extract_plain_text(path)
    if extension == ".md":
        return _extract_markdown(path)
    if extension == ".ipynb":
        return _extract_notebook(path)
    if extension == ".py":
        return _extract_python(path)
    if extension in {".r", ".cpp", ".h", ".m"}:
        return _extract_other_code(path, extension)
    if extension == ".vtt":
        return _extract_vtt(path)
    raise ExtractionFailure(
        f"unsupported text extraction extension: {extension or '<none>'}",
        extractor_name=_extractor_name_for_extension(extension),
        extractor_version=None,
        metadata_json=_json_dumps({"extension": extension}),
    )


def _extract_pdf(path: Path, config: Config) -> tuple[RawChunk, ...]:
    import fitz

    raw_chunks: list[RawChunk] = []
    page_count = 0
    with fitz.open(path) as document:
        page_count = document.page_count
        for page_index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                raw_chunks.append(
                    RawChunk(
                        source_type="document",
                        title=_title_from_text(text),
                        text=text,
                        location_type="page",
                        location_value=str(page_index),
                        metadata={"page": page_index},
                    )
                )

        total_text = "\n".join(chunk.text for chunk in raw_chunks)
        if page_count and len(total_text.strip()) < PDF_SCANNED_TEXT_CHAR_THRESHOLD:
            if not config.ocr_enabled:
                raise ExtractionFailure(
                    SCANNED_PDF_OCR_REASON,
                    extractor_name="pdf-pymupdf",
                    extractor_version=_package_version("PyMuPDF"),
                    metadata_json=_json_dumps(
                        {"page_count": page_count, "ocr_enabled": False}
                    ),
                )
            raw_chunks = _ocr_pdf_pages(document)

    return tuple(raw_chunks)


def _ocr_pdf_pages(document: Any) -> list[RawChunk]:
    if shutil.which("tesseract") is None:
        raise ExtractionFailure(
            SCANNED_PDF_OCR_REASON,
            extractor_name="pdf-pymupdf-ocr",
            extractor_version=_package_version("pytesseract"),
            metadata_json=_json_dumps({"tesseract_available": False}),
        )

    import fitz
    import pytesseract
    from PIL import Image

    try:
        pytesseract.get_tesseract_version()
    except pytesseract.TesseractNotFoundError as exc:
        raise ExtractionFailure(
            SCANNED_PDF_OCR_REASON,
            extractor_name="pdf-pymupdf-ocr",
            extractor_version=_package_version("pytesseract"),
            metadata_json=_json_dumps({"tesseract_available": False}),
        ) from exc

    raw_chunks: list[RawChunk] = []
    matrix = fitz.Matrix(2, 2)
    for page_index, page in enumerate(document, start=1):
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))
        text = pytesseract.image_to_string(image).strip()
        if text:
            raw_chunks.append(
                RawChunk(
                    source_type="document",
                    title=_title_from_text(text),
                    text=text,
                    location_type="page",
                    location_value=str(page_index),
                    metadata={"page": page_index, "ocr": True},
                )
            )
    return raw_chunks


def _extract_pptx(path: Path) -> tuple[RawChunk, ...]:
    from pptx import Presentation

    presentation = Presentation(path)
    raw_chunks: list[RawChunk] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                text_parts.append(text.strip())

        notes_text = _pptx_notes_text(slide)
        if notes_text:
            text_parts.append(f"Speaker notes:\n{notes_text}")

        text = "\n\n".join(text_parts).strip()
        if not text:
            continue
        title = None
        if slide.shapes.title is not None:
            title_text = getattr(slide.shapes.title, "text", "")
            title = title_text.strip() or None
        raw_chunks.append(
            RawChunk(
                source_type="slides",
                title=title or _title_from_text(text),
                text=text,
                location_type="slide",
                location_value=str(slide_index),
                metadata={"slide": slide_index, "has_speaker_notes": bool(notes_text)},
            )
        )
    return tuple(raw_chunks)


def _pptx_notes_text(slide: Any) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
    except (AttributeError, KeyError, ValueError):
        return ""
    paragraphs = [
        paragraph.text.strip()
        for paragraph in notes_frame.paragraphs
        if paragraph.text and paragraph.text.strip()
    ]
    return "\n".join(paragraphs).strip()


def _extract_docx(path: Path) -> tuple[RawChunk, ...]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(path)
    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            text = Paragraph(child, document).text.strip()
            if text:
                blocks.append(text)
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(row for row in rows if row.strip()).strip()
            if table_text:
                blocks.append(table_text)

    return tuple(
        _group_text_blocks(
            blocks,
            source_type="document",
            location_type="docx_section",
            title_prefix="DOCX section",
            metadata={"format": "docx"},
        )
    )


def _extract_plain_text(path: Path) -> tuple[RawChunk, ...]:
    text = _read_text_file(path)
    blocks = _paragraph_blocks(text)
    return tuple(
        _group_text_blocks(
            blocks,
            source_type="document",
            location_type="text_section",
            title_prefix="Text section",
            metadata={"format": "text"},
        )
    )


def _extract_markdown(path: Path) -> tuple[RawChunk, ...]:
    text = _read_text_file(path)
    sections = _markdown_sections(text)
    raw_chunks: list[RawChunk] = []
    for section_index, (title, section_text) in enumerate(sections, start=1):
        if not section_text.strip():
            continue
        raw_chunks.append(
            RawChunk(
                source_type="document",
                title=title or _title_from_text(section_text),
                text=section_text.strip(),
                location_type="markdown_section",
                location_value=str(section_index),
                metadata={"section": section_index, "format": "markdown"},
            )
        )
    if raw_chunks:
        return tuple(raw_chunks)
    return _extract_plain_text(path)


def _extract_notebook(path: Path) -> tuple[RawChunk, ...]:
    import nbformat

    notebook = nbformat.read(path, as_version=4)
    raw_chunks: list[RawChunk] = []
    for cell_index, cell in enumerate(notebook.cells, start=1):
        cell_type = cell.get("cell_type")
        if cell_type not in {"markdown", "code"}:
            continue
        source = str(cell.get("source", "")).strip()
        text = source
        output_text = ""
        if cell_type == "code":
            output_text = _notebook_output_text(cell.get("outputs", []))
            if output_text:
                text = f"{source}\n\nOutput:\n{output_text}".strip()
        if not text:
            continue
        raw_chunks.append(
            RawChunk(
                source_type="notebook",
                title=f"{cell_type} cell {cell_index}",
                text=text,
                location_type="notebook_cell",
                location_value=str(cell_index),
                metadata={
                    "cell": cell_index,
                    "cell_type": cell_type,
                    "output_truncated_to_chars": (
                        NOTEBOOK_OUTPUT_CHAR_LIMIT if output_text else None
                    ),
                },
            )
        )
    return tuple(raw_chunks)


def _extract_python(path: Path) -> tuple[RawChunk, ...]:
    source = _read_text_file(path)
    source_lines = source.splitlines()
    tree = ast.parse(source, filename=str(path))
    raw_chunks: list[RawChunk] = []
    handled_line_numbers: set[int] = set()

    import_nodes = [
        node for node in tree.body if isinstance(node, (ast.Import, ast.ImportFrom))
    ]
    if import_nodes:
        import_lines = _source_for_nodes(source_lines, import_nodes)
        handled_line_numbers.update(_line_numbers_for_nodes(import_nodes))
        raw_chunks.append(
            RawChunk(
                source_type="code",
                title="Imports",
                text=import_lines,
                location_type="module",
                location_value="imports",
                metadata={"language": "python", "section": "imports"},
            )
        )

    module_docstring = ast.get_docstring(tree)
    if module_docstring:
        module_docstring_node = _module_docstring_node(tree)
        if module_docstring_node is not None:
            handled_line_numbers.update(
                _line_numbers_for_nodes([module_docstring_node])
            )
        raw_chunks.append(
            RawChunk(
                source_type="code",
                title="Module docstring",
                text=module_docstring,
                location_type="module",
                location_value="docstring",
                metadata={"language": "python", "section": "docstring"},
            )
        )

    for node in tree.body:
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
            handled_line_numbers.update(_line_numbers_for_nodes([node]))
            node_text = _source_for_node(source_lines, node)
            location_type = "class" if isinstance(node, ast.ClassDef) else "function"
            raw_chunks.append(
                RawChunk(
                    source_type="code",
                    title=node.name,
                    text=node_text,
                    location_type=location_type,
                    location_value=node.name,
                    metadata={
                        "language": "python",
                        "line_start": node.lineno,
                        "line_end": getattr(node, "end_lineno", node.lineno),
                    },
                )
            )

    module_lines = [
        line
        for line_number, line in enumerate(source_lines, start=1)
        if line_number not in handled_line_numbers and line.strip()
    ]
    module_text = "\n".join(module_lines).strip()
    if module_text:
        raw_chunks.append(
            RawChunk(
                source_type="code",
                title="Module",
                text=module_text,
                location_type="module",
                location_value="module",
                metadata={"language": "python", "section": "module"},
            )
        )

    return tuple(raw_chunks)


def _extract_other_code(path: Path, extension: str) -> tuple[RawChunk, ...]:
    text = _read_text_file(path)
    language = {
        ".r": "r",
        ".cpp": "cpp",
        ".h": "cpp-header",
        ".m": "matlab",
    }[extension]
    matches = _function_matches_for_code(text, extension)
    if not matches:
        return (
            RawChunk(
                source_type="code",
                title=path.name,
                text=text.strip(),
                location_type="module",
                location_value="module",
                metadata={"language": language, "fallback": "whole_file"},
            ),
        )

    raw_chunks: list[RawChunk] = []
    for index, (name, start, line_number) in enumerate(matches):
        end = matches[index + 1][1] if index + 1 < len(matches) else len(text)
        chunk_text = text[start:end].strip()
        if not chunk_text:
            continue
        raw_chunks.append(
            RawChunk(
                source_type="code",
                title=name,
                text=chunk_text,
                location_type="function",
                location_value=name,
                metadata={
                    "language": language,
                    "line_start": line_number,
                    "fallback": "regex",
                },
            )
        )
    return tuple(raw_chunks)


def _extract_vtt(path: Path) -> tuple[RawChunk, ...]:
    text = _read_text_file(path)
    lines = text.splitlines()
    raw_chunks: list[RawChunk] = []
    index = 0
    while index < len(lines):
        line = lines[index].strip()
        match = VTT_TIMESTAMP_RE.match(line)
        if not match:
            index += 1
            continue
        start = match.group("start")
        end = match.group("end")
        index += 1
        cue_lines: list[str] = []
        while index < len(lines) and lines[index].strip():
            cue_lines.append(lines[index].strip())
            index += 1
        cue_text = "\n".join(cue_lines).strip()
        if cue_text:
            raw_chunks.append(
                RawChunk(
                    source_type="transcript",
                    title=start,
                    text=cue_text,
                    location_type="timestamp",
                    location_value=start,
                    metadata={"timestamp_start": start, "timestamp_end": end},
                )
            )
    return tuple(raw_chunks)


def _finalize_chunks(
    *,
    file_record: PendingFileRecord,
    raw_chunks: Iterable[RawChunk],
    max_tokens: int,
) -> tuple[ChunkRecord, ...]:
    chunks: list[ChunkRecord] = []
    for raw_chunk in raw_chunks:
        text = raw_chunk.text.strip()
        if not text:
            continue
        pieces = _split_text_by_tokens(text, max_tokens)
        for piece_index, piece in enumerate(pieces, start=1):
            is_subchunk = len(pieces) > 1
            metadata_payload = dict(raw_chunk.metadata)
            metadata_payload.update(
                {
                    "source_location_type": raw_chunk.location_type,
                    "source_location_value": raw_chunk.location_value,
                }
            )
            if is_subchunk:
                metadata_payload.update(
                    {
                        "subchunk_index": piece_index,
                        "subchunk_count": len(pieces),
                    }
                )
                location_type = "subchunk"
                location_value = (
                    f"{raw_chunk.location_type}:{raw_chunk.location_value}:"
                    f"part:{piece_index}"
                )
            else:
                location_type = raw_chunk.location_type
                location_value = raw_chunk.location_value

            chunk_index = len(chunks)
            chunks.append(
                ChunkRecord(
                    file_id=file_record.id,
                    chunk_uid=f"file-{file_record.id}-chunk-{chunk_index}",
                    source_type=raw_chunk.source_type,
                    chunk_index=chunk_index,
                    title=raw_chunk.title,
                    text=piece,
                    token_count=_count_tokens(piece),
                    location_type=location_type,
                    location_value=location_value,
                    metadata_json=_json_dumps(metadata_payload),
                )
            )
    return tuple(chunks)


def _record_failed_extraction(
    connection: sqlite3.Connection,
    *,
    run_id: int,
    file_record: PendingFileRecord,
    extractor_name: str,
    extractor_version: str | None,
    metadata_json: str | None,
    error: str,
) -> None:
    with connection:
        _persist_failed_extraction(
            connection,
            run_id=run_id,
            file_record=file_record,
            extractor_name=extractor_name,
            extractor_version=extractor_version,
            metadata_json=metadata_json,
            error=error,
            extracted_at=_utc_now(),
        )


def _persist_successful_extraction(
    connection: sqlite3.Connection,
    *,
    run_id: int,
    extracted: ExtractedDocument,
    extracted_at: str,
) -> None:
    _delete_existing_chunks(connection, extracted.file_id)
    extracted_document_id = _upsert_extracted_document(
        connection,
        run_id=run_id,
        file_id=extracted.file_id,
        extractor_name=extracted.extractor_name,
        extractor_version=extracted.extractor_version,
        status=extracted.status,
        text_length=extracted.text_length,
        chunk_count=extracted.chunk_count,
        metadata_json=extracted.metadata_json,
        error=None,
        extracted_at=extracted_at,
    )
    for chunk in extracted.chunks:
        connection.execute(
            """
            INSERT INTO chunks (
                file_id,
                extracted_document_id,
                chunk_uid,
                source_type,
                chunk_index,
                title,
                text,
                token_count,
                location_type,
                location_value,
                metadata_json,
                created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                chunk.file_id,
                extracted_document_id,
                chunk.chunk_uid,
                chunk.source_type,
                chunk.chunk_index,
                chunk.title,
                chunk.text,
                chunk.token_count,
                chunk.location_type,
                chunk.location_value,
                chunk.metadata_json,
                extracted_at,
            ),
        )
    connection.execute(
        """
        UPDATE files
        SET index_status = 'indexed',
            reason_not_indexed = NULL
        WHERE id = ?
        """,
        (extracted.file_id,),
    )


def _persist_failed_extraction(
    connection: sqlite3.Connection,
    *,
    run_id: int,
    file_record: PendingFileRecord,
    extractor_name: str,
    extractor_version: str | None,
    metadata_json: str | None,
    error: str,
    extracted_at: str,
) -> None:
    _delete_existing_chunks(connection, file_record.id)
    _upsert_extracted_document(
        connection,
        run_id=run_id,
        file_id=file_record.id,
        extractor_name=extractor_name,
        extractor_version=extractor_version,
        status="failed",
        text_length=0,
        chunk_count=0,
        metadata_json=metadata_json or _json_dumps({"extension": file_record.extension}),
        error=error,
        extracted_at=extracted_at,
    )
    connection.execute(
        """
        UPDATE files
        SET index_status = 'failed',
            reason_not_indexed = ?
        WHERE id = ?
        """,
        (_truncate(error, ERROR_CHAR_LIMIT), file_record.id),
    )


def _delete_existing_chunks(connection: sqlite3.Connection, file_id: int) -> None:
    connection.execute(
        """
        DELETE FROM embeddings
        WHERE chunk_id IN (
            SELECT id
            FROM chunks
            WHERE file_id = ?
        )
        """,
        (file_id,),
    )
    connection.execute(
        """
        DELETE FROM chunk_fts
        WHERE chunk_id IN (
            SELECT id
            FROM chunks
            WHERE file_id = ?
        )
        """,
        (file_id,),
    )
    connection.execute("DELETE FROM chunks WHERE file_id = ?", (file_id,))


def _upsert_extracted_document(
    connection: sqlite3.Connection,
    *,
    run_id: int,
    file_id: int,
    extractor_name: str,
    extractor_version: str | None,
    status: str,
    text_length: int,
    chunk_count: int,
    metadata_json: str,
    error: str | None,
    extracted_at: str,
) -> int:
    connection.execute(
        """
        INSERT INTO extracted_documents (
            file_id,
            extraction_run_id,
            extractor_name,
            extractor_version,
            status,
            text_length,
            chunk_count,
            metadata_json,
            error,
            extracted_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(file_id, extractor_name) DO UPDATE SET
            extraction_run_id = excluded.extraction_run_id,
            extractor_version = excluded.extractor_version,
            status = excluded.status,
            text_length = excluded.text_length,
            chunk_count = excluded.chunk_count,
            metadata_json = excluded.metadata_json,
            error = excluded.error,
            extracted_at = excluded.extracted_at
        """,
        (
            file_id,
            run_id,
            extractor_name,
            extractor_version,
            status,
            text_length,
            chunk_count,
            metadata_json,
            error,
            extracted_at,
        ),
    )
    row = connection.execute(
        """
        SELECT id
        FROM extracted_documents
        WHERE file_id = ? AND extractor_name = ?
        """,
        (file_id, extractor_name),
    ).fetchone()
    if row is None:
        raise ExtractionError(f"Failed to upsert extracted document for file {file_id}")
    return int(row["id"])


def _load_pending_files(
    connection: sqlite3.Connection,
    category: str | None,
) -> tuple[PendingFileRecord, ...]:
    categories = [category] if category is not None else sorted(TEXT_EXTRACTION_CATEGORIES)
    placeholders = ",".join("?" for _ in categories)
    rows = connection.execute(
        f"""
        SELECT id, path, relative_path, filename, extension, category, content_hash
        FROM files
        WHERE index_status = 'pending'
          AND category IN ({placeholders})
        ORDER BY relative_path COLLATE NOCASE
        """,
        categories,
    ).fetchall()
    return tuple(
        PendingFileRecord(
            id=int(row["id"]),
            path=Path(str(row["path"])),
            relative_path=str(row["relative_path"]),
            filename=str(row["filename"]),
            extension=str(row["extension"]),
            category=str(row["category"]),
            content_hash=row["content_hash"],
        )
        for row in rows
    )


def _start_extraction_run(
    connection: sqlite3.Connection,
    config: Config,
    category: str | None,
    started_at: str,
) -> int:
    config_json = _json_dumps(
        {
            "run_type": "extraction",
            "category": category,
            "courses_root": str(config.courses_root),
            "data_dir": str(config.data_dir),
            "sqlite_path": str(config.sqlite_path),
            "ocr_enabled": config.ocr_enabled,
            "max_chunk_tokens": DEFAULT_MAX_CHUNK_TOKENS,
            "handled_categories": sorted(TEXT_EXTRACTION_CATEGORIES),
        }
    )
    cursor = connection.execute(
        """
        INSERT INTO extraction_runs (
            started_at,
            status,
            config_json,
            files_seen,
            files_indexed,
            files_metadata_only,
            files_failed
        )
        VALUES (?, 'running', ?, 0, 0, 0, 0)
        """,
        (started_at, config_json),
    )
    return int(cursor.lastrowid)


def _finish_extraction_run(
    connection: sqlite3.Connection,
    *,
    run_id: int,
    finished_at: str,
    status: str,
    files_seen: int,
    files_indexed: int,
    files_failed: int,
    error: str | None,
) -> None:
    connection.execute(
        """
        UPDATE extraction_runs
        SET finished_at = ?,
            status = ?,
            files_seen = ?,
            files_indexed = ?,
            files_metadata_only = 0,
            files_failed = ?,
            error = ?
        WHERE id = ?
        """,
        (
            finished_at,
            status,
            files_seen,
            files_indexed,
            files_failed,
            error,
            run_id,
        ),
    )


def _latest_extraction_run(
    connection: sqlite3.Connection,
) -> tuple[int | None, str | None]:
    rows = connection.execute(
        """
        SELECT id, started_at, config_json
        FROM extraction_runs
        ORDER BY id DESC
        """
    ).fetchall()
    for row in rows:
        try:
            payload = json.loads(row["config_json"])
        except json.JSONDecodeError:
            continue
        if payload.get("run_type") == "extraction":
            return int(row["id"]), str(row["started_at"])
    return None, None


def _count_pending_by_category(connection: sqlite3.Connection) -> dict[str, int]:
    placeholders = ",".join("?" for _ in TEXT_EXTRACTION_CATEGORIES)
    rows = connection.execute(
        f"""
        SELECT category, COUNT(*) AS count
        FROM files
        WHERE index_status = 'pending'
          AND category IN ({placeholders})
        GROUP BY category
        ORDER BY category
        """,
        sorted(TEXT_EXTRACTION_CATEGORIES),
    ).fetchall()
    return {str(row["category"]): int(row["count"]) for row in rows}


def _count_chunks_by_source_type(connection: sqlite3.Connection) -> dict[str, int]:
    rows = connection.execute(
        """
        SELECT source_type, COUNT(*) AS count
        FROM chunks
        GROUP BY source_type
        ORDER BY source_type
        """
    ).fetchall()
    return {str(row["source_type"]): int(row["count"]) for row in rows}


def _count_files_by_status(connection: sqlite3.Connection, status: str) -> int:
    placeholders = ",".join("?" for _ in TEXT_EXTRACTION_CATEGORIES)
    row = connection.execute(
        f"""
        SELECT COUNT(*) AS count
        FROM files
        WHERE index_status = ?
          AND category IN ({placeholders})
        """,
        (status, *sorted(TEXT_EXTRACTION_CATEGORIES)),
    ).fetchone()
    return int(row["count"] if row else 0)


def _count_table_rows(connection: sqlite3.Connection, table: str) -> int:
    row = connection.execute(f"SELECT COUNT(*) AS count FROM {table}").fetchone()
    return int(row["count"] if row else 0)


def _load_recent_failures(
    connection: sqlite3.Connection,
) -> tuple[ExtractionFailureSummary, ...]:
    rows = connection.execute(
        """
        SELECT files.id AS file_id, files.path AS path, extracted_documents.error AS error
        FROM extracted_documents
        JOIN files ON files.id = extracted_documents.file_id
        WHERE extracted_documents.status = 'failed'
        ORDER BY extracted_documents.extracted_at DESC
        LIMIT 10
        """
    ).fetchall()
    return tuple(
        ExtractionFailureSummary(
            file_id=int(row["file_id"]),
            path=str(row["path"]),
            error=str(row["error"]),
        )
        for row in rows
    )


def _validate_category(category: str | None) -> None:
    if category is None:
        return
    if category not in TEXT_EXTRACTION_CATEGORIES:
        allowed = ", ".join(sorted(TEXT_EXTRACTION_CATEGORIES))
        raise ExtractionError(
            f"Text extraction category must be one of: {allowed}. "
            "Data schema summaries are handled by Feature 05."
        )


def _failure_from_exception(
    file_record: PendingFileRecord,
    exc: ExtractionFailure,
) -> ExtractionFailureSummary:
    return ExtractionFailureSummary(
        file_id=file_record.id,
        path=str(file_record.path),
        error=_truncate(str(exc), ERROR_CHAR_LIMIT),
    )


def _read_text_file(path: Path) -> str:
    last_error: UnicodeDecodeError | None = None
    for encoding in TEXT_ENCODINGS:
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError as exc:
            last_error = exc
    if last_error is not None:
        raise last_error
    return path.read_text()


def _paragraph_blocks(text: str) -> list[str]:
    blocks = re.split(r"\n\s*\n", text)
    return [block.strip() for block in blocks if block.strip()]


def _markdown_sections(text: str) -> list[tuple[str | None, str]]:
    sections: list[tuple[str | None, str]] = []
    current_title: str | None = None
    current_lines: list[str] = []
    for line in text.splitlines():
        heading = re.match(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$", line)
        if heading and current_lines:
            sections.append((current_title, "\n".join(current_lines).strip()))
            current_lines = []
        if heading:
            current_title = heading.group(2).strip()
        current_lines.append(line)
    if current_lines:
        sections.append((current_title, "\n".join(current_lines).strip()))
    return sections


def _group_text_blocks(
    blocks: Iterable[str],
    *,
    source_type: str,
    location_type: str,
    title_prefix: str,
    metadata: Mapping[str, object],
) -> list[RawChunk]:
    raw_chunks: list[RawChunk] = []
    current: list[str] = []
    current_tokens = 0

    for block in blocks:
        block = block.strip()
        if not block:
            continue
        block_tokens = _count_tokens(block)
        if current and current_tokens + block_tokens > DEFAULT_MAX_CHUNK_TOKENS:
            section_index = len(raw_chunks) + 1
            text = "\n\n".join(current).strip()
            raw_chunks.append(
                RawChunk(
                    source_type=source_type,
                    title=f"{title_prefix} {section_index}",
                    text=text,
                    location_type=location_type,
                    location_value=str(section_index),
                    metadata={**metadata, "section": section_index},
                )
            )
            current = []
            current_tokens = 0
        current.append(block)
        current_tokens += block_tokens

    if current:
        section_index = len(raw_chunks) + 1
        text = "\n\n".join(current).strip()
        raw_chunks.append(
            RawChunk(
                source_type=source_type,
                title=f"{title_prefix} {section_index}",
                text=text,
                location_type=location_type,
                location_value=str(section_index),
                metadata={**metadata, "section": section_index},
            )
        )
    return raw_chunks


def _notebook_output_text(outputs: Iterable[Mapping[str, Any]]) -> str:
    pieces: list[str] = []
    for output in outputs:
        output_type = output.get("output_type")
        if output_type == "stream":
            text = output.get("text", "")
            pieces.append(_stringify_notebook_text(text))
        elif output_type in {"execute_result", "display_data"}:
            data = output.get("data", {})
            if isinstance(data, Mapping) and "text/plain" in data:
                pieces.append(_stringify_notebook_text(data["text/plain"]))
        elif output_type == "error":
            trace_lines = output.get("traceback", [])
            pieces.append(_stringify_notebook_text(trace_lines))

    text = "\n".join(piece for piece in pieces if piece.strip()).strip()
    return _truncate(text, NOTEBOOK_OUTPUT_CHAR_LIMIT)


def _stringify_notebook_text(value: object) -> str:
    if isinstance(value, list):
        return "".join(str(item) for item in value)
    return str(value)


def _source_for_nodes(source_lines: list[str], nodes: Iterable[ast.AST]) -> str:
    return "\n".join(_source_for_node(source_lines, node) for node in nodes).strip()


def _source_for_node(source_lines: list[str], node: ast.AST) -> str:
    start = getattr(node, "lineno", 1)
    end = getattr(node, "end_lineno", start)
    return "\n".join(source_lines[start - 1 : end]).strip()


def _module_docstring_node(tree: ast.Module) -> ast.Expr | None:
    if not tree.body:
        return None
    node = tree.body[0]
    if isinstance(node, ast.Expr) and isinstance(node.value, ast.Constant):
        if isinstance(node.value.value, str):
            return node
    return None


def _line_numbers_for_nodes(nodes: Iterable[ast.AST]) -> set[int]:
    line_numbers: set[int] = set()
    for node in nodes:
        start = getattr(node, "lineno", 1)
        end = getattr(node, "end_lineno", start)
        line_numbers.update(range(start, end + 1))
    return line_numbers


def _function_matches_for_code(
    text: str,
    extension: str,
) -> list[tuple[str, int, int]]:
    if extension == ".r":
        regex = R_FUNCTION_RE
    elif extension in {".cpp", ".h"}:
        regex = CPP_FUNCTION_RE
    else:
        regex = MATLAB_FUNCTION_RE

    matches: list[tuple[str, int, int]] = []
    for match in regex.finditer(text):
        name = match.groupdict().get("name") or match.group(1)
        line_number = text.count("\n", 0, match.start()) + 1
        matches.append((name, match.start(), line_number))
    return matches


def _split_text_by_tokens(text: str, max_tokens: int) -> list[str]:
    words = text.split()
    if len(words) <= max_tokens:
        return [text]
    return [
        " ".join(words[index : index + max_tokens]).strip()
        for index in range(0, len(words), max_tokens)
    ]


def _count_tokens(text: str) -> int:
    return len(text.split())


def _title_from_text(text: str, max_length: int = 80) -> str | None:
    for line in text.splitlines():
        normalized = " ".join(line.split())
        if normalized:
            return _truncate(normalized, max_length)
    return None


def _extractor_name_for_extension(extension: str) -> str:
    return {
        ".pdf": "pdf-pymupdf",
        ".pptx": "pptx-python-pptx",
        ".ppt": "legacy-ppt-unsupported",
        ".docx": "docx-python-docx",
        ".doc": "legacy-doc-unsupported",
        ".txt": "plain-text",
        ".md": "markdown-text",
        ".ipynb": "notebook-nbformat",
        ".py": "python-ast",
        ".r": "code-regex",
        ".cpp": "code-regex",
        ".h": "code-regex",
        ".m": "code-regex",
        ".vtt": "vtt-parser",
    }.get(extension, "unsupported-extractor")


def _extractor_version_for_extension(extension: str) -> str | None:
    package_name = {
        ".pdf": "PyMuPDF",
        ".pptx": "python-pptx",
        ".docx": "python-docx",
        ".ipynb": "nbformat",
        ".doc": None,
        ".ppt": None,
    }.get(extension)
    if package_name is None:
        return None
    return _package_version(package_name)


def _package_version(package_name: str) -> str | None:
    try:
        return metadata.version(package_name)
    except metadata.PackageNotFoundError:
        return None


def _json_dumps(payload: Mapping[str, object]) -> str:
    return json.dumps(payload, sort_keys=True)


def _format_exception(exc: Exception) -> str:
    message = "".join(traceback.format_exception_only(type(exc), exc)).strip()
    return _truncate(message, ERROR_CHAR_LIMIT)


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return f"{text[: max_chars - 3]}..."


def _utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()
