"""PowerPoint slide extraction."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from .._textutils import _title_from_text
from ..models import RawChunk


def _extract_pptx(path: Path) -> tuple[RawChunk, ...]:
    from pptx import Presentation

    presentation = Presentation(path)
    raw_chunks: list[RawChunk] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                text_parts.append(text.strip())

        notes_text = _pptx_notes_text(slide)
        if notes_text:
            text_parts.append(f"Speaker notes:\n{notes_text}")

        text = "\n\n".join(text_parts).strip()
        if not text:
            continue
        title = None
        if slide.shapes.title is not None:
            title_text = getattr(slide.shapes.title, "text", "")
            title = title_text.strip() or None
        raw_chunks.append(
            RawChunk(
                source_type="slides",
                title=title or _title_from_text(text),
                text=text,
                location_type="slide",
                location_value=str(slide_index),
                metadata={"slide": slide_index, "has_speaker_notes": bool(notes_text)},
            )
        )
    return tuple(raw_chunks)


def _pptx_notes_text(slide: Any) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
    except (AttributeError, KeyError, ValueError):
        return ""
    paragraphs = [
        paragraph.text.strip()
        for paragraph in notes_frame.paragraphs
        if paragraph.text and paragraph.text.strip()
    ]
    return "\n".join(paragraphs).strip()
