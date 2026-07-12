from __future__ import annotations

import json
import os
from contextlib import closing
from pathlib import Path

import nbformat
import pytest

from uni_rag_agent.extraction import (
    DEFAULT_MAX_CHUNK_TOKENS,
    LEGACY_FORMAT_REASON,
    ExtractionFailure,
    PendingFileRecord,
    SCANNED_PDF_OCR_REASON,
    extract_file,
    extract_pending_files,
    load_extraction_status,
)
from uni_rag_agent.extraction.constants import NO_TEXT_REASON
from uni_rag_agent.inventory import inventory_courses
from uni_rag_agent.storage import connect_sqlite
from tests.sqlite_helpers import insert_search_result
from tests.support import make_config

REPO_ROOT = Path(__file__).resolve().parents[1]


def _pending_file(
    path: Path,
    *,
    file_id: int,
    category: str,
    content_hash: str | None = "test-hash",
) -> PendingFileRecord:
    return PendingFileRecord(
        id=file_id,
        path=path,
        relative_path=path.name,
        filename=path.name,
        extension=path.suffix.lower(),
        category=category,
        content_hash=content_hash,
    )


def test_extract_run_processes_text_formats_and_preserves_locations(
    tmp_path: Path,
) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    sentinel = course_dir / "executed.txt"

    (course_dir / "notes.md").write_text(
        "# MapReduce\n\nMapReduce combines map tasks and reduce tasks.",
        encoding="utf-8",
    )
    (course_dir / "syllabus.txt").write_text(
        "BM25 is a keyword ranking method.",
        encoding="utf-8",
    )
    (course_dir / "assignment.py").write_text(
        "\n".join(
            [
                "from pathlib import Path",
                "",
                "Path('executed.txt').write_text('bad')",
                "",
                "class Ranker:",
                "    def score(self):",
                "        return 1",
                "",
                "def train_model():",
                "    return 'trained'",
            ]
        ),
        encoding="utf-8",
    )
    (course_dir / "analysis.r").write_text(
        "fit_model <- function(x) {\n  x + 1\n}",
        encoding="utf-8",
    )
    (course_dir / "kernel.cpp").write_text(
        "int add(int a, int b) {\n  return a + b;\n}",
        encoding="utf-8",
    )
    (course_dir / "script.m").write_text(
        "function y = normalize(x)\ny = x;\nend",
        encoding="utf-8",
    )
    (course_dir / "captions.vtt").write_text(
        "WEBVTT\n\n00:00:01.000 --> 00:00:03.000\nVector search overview.\n",
        encoding="utf-8",
    )
    _write_pdf(course_dir / "paper.pdf", "PDF page text about inverted indexes.")
    _write_docx(course_dir / "report.docx")
    _write_pptx(course_dir / "slides.pptx")
    _write_notebook(course_dir / "search_demo.ipynb")
    (course_dir / "dataset.csv").write_text("term,score\nbm25,1\n", encoding="utf-8")

    inventory_courses(config)
    result = extract_pending_files(config)

    assert result.status == "completed", "multi-format extraction run status"
    assert result.files_seen == 11, "multi-format files seen"
    assert result.files_indexed == 11, "multi-format files indexed"
    assert result.files_failed == 0, "multi-format files failed"
    assert result.by_source_type["document"] >= 4, "document chunk count"
    assert result.by_source_type["slides"] == 1, "slide chunk count"
    assert result.by_source_type["notebook"] == 2, "notebook cell chunk count"
    assert result.by_source_type["code"] >= 6, "code chunk count"
    assert result.by_source_type["transcript"] == 1, "transcript chunk count"
    assert not sentinel.exists(), "python extractor must not execute source code"

    with closing(connect_sqlite(config)) as connection:
        rows = connection.execute(
            """
            SELECT files.filename, files.index_status, files.reason_not_indexed,
                   chunks.source_type, chunks.title, chunks.text,
                   chunks.location_type, chunks.location_value, chunks.metadata_json
            FROM files
            LEFT JOIN chunks ON chunks.file_id = files.id
            ORDER BY files.filename, chunks.chunk_index
            """
        ).fetchall()
        run_row = connection.execute(
            """
            SELECT config_json, files_seen, files_indexed, files_failed
            FROM extraction_runs
            WHERE id = ?
            """,
            (result.run_id,),
        ).fetchone()

    rows_by_file: dict[str, list] = {}
    for row in rows:
        rows_by_file.setdefault(row["filename"], []).append(row)

    assert rows_by_file["dataset.csv"][0]["index_status"] == "pending", (
        "data_schema file should remain pending for Feature 05"
    )
    assert rows_by_file["paper.pdf"][0]["location_type"] == "page", "pdf page"
    assert rows_by_file["paper.pdf"][0]["location_value"] == "1", "pdf page number"
    assert rows_by_file["slides.pptx"][0]["source_type"] == "slides", "pptx source"
    assert rows_by_file["slides.pptx"][0]["location_type"] == "slide", "pptx slide"
    assert rows_by_file["report.docx"][0]["location_type"] == "docx_section", (
        "docx section"
    )
    assert rows_by_file["notes.md"][0]["location_type"] == "markdown_section", (
        "markdown section"
    )
    assert rows_by_file["syllabus.txt"][0]["location_type"] == "text_section", (
        "plain-text section"
    )
    assert rows_by_file["captions.vtt"][0]["location_type"] == "timestamp", "vtt cue"
    assert rows_by_file["captions.vtt"][0]["location_value"] == "00:00:01.000", (
        "vtt cue timestamp"
    )

    notebook_text = "\n".join(row["text"] for row in rows_by_file["search_demo.ipynb"])
    assert "accuracy: 0.95" in notebook_text, "notebook text output"
    assert "image/png" not in notebook_text, "notebook binary output skipped"

    python_locations = {
        (row["location_type"], row["location_value"])
        for row in rows_by_file["assignment.py"]
    }
    assert ("class", "Ranker") in python_locations, "python class location"
    assert ("function", "train_model") in python_locations, "python function location"

    code_titles = {
        row["title"]
        for filename in ("analysis.r", "kernel.cpp", "script.m")
        for row in rows_by_file[filename]
    }
    assert {"fit_model", "add", "normalize"}.issubset(code_titles), (
        "regex code extractor titles"
    )

    run_payload = json.loads(run_row["config_json"])
    assert run_payload["run_type"] == "extraction", "extraction run payload type"
    assert run_row["files_seen"] == 11, "stored extraction files seen"
    assert run_row["files_indexed"] == 11, "stored extraction files indexed"
    assert run_row["files_failed"] == 0, "stored extraction files failed"


def test_extract_file_zero_page_pdf_reports_no_text_reason(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    path = config.courses_root / "empty.pdf"
    _write_zero_page_pdf(path)

    with pytest.raises(ExtractionFailure, match=NO_TEXT_REASON):
        extract_file(_pending_file(path, file_id=1, category="document"), config)


def test_extract_file_docx_table_text_is_preserved(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    path = config.courses_root / "report.docx"
    _write_docx(path)

    extracted = extract_file(
        _pending_file(path, file_id=1, category="document"), config
    )
    text = "\n".join(chunk.text for chunk in extracted.chunks)

    assert "Document paragraph about cosine similarity." in text, "docx paragraph"
    assert "term | weight" in text, "docx table row"
    assert extracted.chunks[0].location_type == "docx_section", "docx location"


def test_extract_file_pptx_speaker_notes_are_included(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    path = config.courses_root / "slides.pptx"
    _write_pptx_with_notes(path)

    extracted = extract_file(_pending_file(path, file_id=1, category="slides"), config)
    chunk = extracted.chunks[0]
    metadata = json.loads(chunk.metadata_json)

    assert chunk.source_type == "slides", "pptx source type"
    assert chunk.location_type == "slide", "pptx slide location"
    assert "Speaker notes:" in chunk.text, "pptx speaker notes label"
    assert "Mention reciprocal rank fusion." in chunk.text, "pptx speaker notes text"
    assert metadata["has_speaker_notes"] is True, "pptx notes metadata"


def test_extract_file_empty_notebook_reports_no_text_reason(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    path = config.courses_root / "empty.ipynb"
    notebook = nbformat.v4.new_notebook()
    notebook.cells = []
    nbformat.write(notebook, path)

    with pytest.raises(ExtractionFailure, match=NO_TEXT_REASON):
        extract_file(_pending_file(path, file_id=1, category="notebook"), config)


def test_python_module_docstring_is_not_duplicated_in_module_chunk(
    tmp_path: Path,
) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    (course_dir / "documented.py").write_text(
        "\n".join(
            [
                '"""Module summary for ranking."""',
                "",
                "from pathlib import Path",
                "",
                "VALUE = 1",
                "",
                "def score():",
                "    return VALUE",
            ]
        ),
        encoding="utf-8",
    )

    inventory_courses(config)
    result = extract_pending_files(config, category="code")

    assert result.files_indexed == 1

    with closing(connect_sqlite(config)) as connection:
        rows = connection.execute(
            """
            SELECT chunks.title, chunks.text, chunks.location_value
            FROM chunks
            JOIN files ON files.id = chunks.file_id
            WHERE files.filename = ?
            ORDER BY chunks.chunk_index
            """,
            ("documented.py",),
        ).fetchall()

    rows_by_title = {row["title"]: row for row in rows}
    assert rows_by_title["Module docstring"]["text"] == "Module summary for ranking."
    assert "from pathlib import Path" in rows_by_title["Imports"]["text"]
    assert "def score" in rows_by_title["score"]["text"]

    module_row = next(row for row in rows if row["location_value"] == "module")
    assert "VALUE = 1" in module_row["text"]
    assert "Module summary for ranking" not in module_row["text"]
    assert "def score" not in module_row["text"]


def test_legacy_formats_fail_per_file_with_expected_reason(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    (course_dir / "legacy.doc").write_text("legacy document", encoding="utf-8")
    (course_dir / "legacy.ppt").write_text("legacy slides", encoding="utf-8")
    (course_dir / "syllabus.txt").write_text("BM25", encoding="utf-8")

    inventory_courses(config)
    result = extract_pending_files(config)

    assert result.files_seen == 3
    assert result.files_indexed == 1
    assert result.files_failed == 2
    assert all(LEGACY_FORMAT_REASON in failure.error for failure in result.failures)

    with closing(connect_sqlite(config)) as connection:
        rows = connection.execute(
            """
            SELECT filename, index_status, reason_not_indexed
            FROM files
            ORDER BY filename
            """
        ).fetchall()

    by_name = {row["filename"]: row for row in rows}
    assert by_name["legacy.doc"]["index_status"] == "failed"
    assert by_name["legacy.doc"]["reason_not_indexed"] == LEGACY_FORMAT_REASON
    assert by_name["legacy.ppt"]["index_status"] == "failed"
    assert by_name["legacy.ppt"]["reason_not_indexed"] == LEGACY_FORMAT_REASON
    assert by_name["syllabus.txt"]["index_status"] == "indexed"


def test_file_failure_does_not_abort_extraction_run(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    (course_dir / "broken.pdf").write_bytes(b"not a pdf")
    (course_dir / "notes.md").write_text("# Search\n\nBM25", encoding="utf-8")

    inventory_courses(config)
    result = extract_pending_files(config)

    assert result.status == "completed"
    assert result.files_seen == 2
    assert result.files_indexed == 1
    assert result.files_failed == 1
    assert "broken.pdf" in result.failures[0].path

    with closing(connect_sqlite(config)) as connection:
        rows = connection.execute(
            """
            SELECT filename, index_status
            FROM files
            ORDER BY filename
            """
        ).fetchall()

    by_name = {row["filename"]: row for row in rows}
    assert by_name["broken.pdf"]["index_status"] == "failed"
    assert by_name["notes.md"]["index_status"] == "indexed"


def test_reextraction_nulls_historical_search_result_chunk_reference(
    tmp_path: Path,
) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    source_file = course_dir / "notes.md"
    source_file.write_text("# Search\n\nBM25 first version", encoding="utf-8")

    inventory_courses(config)
    first_result = extract_pending_files(config)

    with closing(connect_sqlite(config)) as connection:
        chunk_row = connection.execute(
            """
            SELECT chunks.id AS chunk_id, files.id AS file_id
            FROM chunks
            JOIN files ON files.id = chunks.file_id
            WHERE files.filename = ?
            """,
            ("notes.md",),
        ).fetchone()
        search_result = insert_search_result(
            connection,
            chunk_id=chunk_row["chunk_id"],
            file_id=chunk_row["file_id"],
            started_at=first_result.started_at,
            finished_at=first_result.finished_at,
        )
        connection.commit()

    source_file.write_text("# Search\n\nBM25 changed version", encoding="utf-8")
    current_mtime = source_file.stat().st_mtime
    os.utime(source_file, (current_mtime + 5, current_mtime + 5))
    inventory_courses(config)
    second_result = extract_pending_files(config)

    assert second_result.status == "completed"
    assert second_result.files_indexed == 1
    assert second_result.files_failed == 0

    with closing(connect_sqlite(config)) as connection:
        historical_result = connection.execute(
            """
            SELECT chunk_id, file_id
            FROM search_results
            WHERE id = ?
            """,
            (search_result.search_result_id,),
        ).fetchone()
        current_chunk_text = connection.execute(
            """
            SELECT chunks.text
            FROM chunks
            JOIN files ON files.id = chunks.file_id
            WHERE files.filename = ?
            """,
            ("notes.md",),
        ).fetchone()

    assert historical_result["chunk_id"] is None
    assert historical_result["file_id"] == chunk_row["file_id"]
    assert "BM25 changed version" in current_chunk_text["text"]


def test_scanned_pdf_without_ocr_uses_contract_reason(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    _write_blank_pdf(course_dir / "scanned.pdf")

    inventory_courses(config)
    result = extract_pending_files(config)

    assert result.files_failed == 1
    assert result.failures[0].error == SCANNED_PDF_OCR_REASON

    with closing(connect_sqlite(config)) as connection:
        row = connection.execute(
            "SELECT index_status, reason_not_indexed FROM files WHERE filename = ?",
            ("scanned.pdf",),
        ).fetchone()

    assert row["index_status"] == "failed"
    assert row["reason_not_indexed"] == SCANNED_PDF_OCR_REASON


def test_overlarge_units_are_subchunked(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    words = [f"term{i}" for i in range(DEFAULT_MAX_CHUNK_TOKENS + 25)]
    (course_dir / "long.txt").write_text(" ".join(words), encoding="utf-8")

    inventory_courses(config)
    result = extract_pending_files(config)

    assert result.files_indexed == 1
    assert result.chunks_created == 2

    with closing(connect_sqlite(config)) as connection:
        rows = connection.execute(
            """
            SELECT location_type, location_value, token_count, metadata_json
            FROM chunks
            ORDER BY chunk_index
            """
        ).fetchall()

    assert [row["location_type"] for row in rows] == ["subchunk", "subchunk"]
    assert rows[0]["token_count"] == DEFAULT_MAX_CHUNK_TOKENS
    assert rows[1]["token_count"] == 25
    metadata_payload = json.loads(rows[0]["metadata_json"])
    assert metadata_payload["source_location_type"] == "text_section"
    assert metadata_payload["subchunk_count"] == 2


def test_extract_run_category_filter_and_status(tmp_path: Path) -> None:
    config = make_config(tmp_path)
    course_dir = config.courses_root / "Information Retrieval"
    course_dir.mkdir()
    (course_dir / "notes.md").write_text("# Search\n\nBM25", encoding="utf-8")
    (course_dir / "assignment.py").write_text(
        "def score():\n    return 1\n", encoding="utf-8"
    )

    inventory_courses(config)
    result = extract_pending_files(config, category="document")
    status = load_extraction_status(config)

    assert result.files_seen == 1
    assert result.files_indexed == 1
    assert status.latest_extraction_run_id == result.run_id
    assert status.indexed_text_files == 1
    assert status.pending_text_files == 1
    assert status.pending_by_category["code"] == 1
    assert status.chunks_by_source_type["document"] == 1


def test_extraction_eda_notebook_is_valid_and_read_only() -> None:
    notebook_path = REPO_ROOT / "notebooks" / "extraction_eda.ipynb"
    notebook = nbformat.read(notebook_path, as_version=4)
    source_text = "\n".join(cell.get("source", "") for cell in notebook.cells)
    cell_ids = {cell.get("id") for cell in notebook.cells}

    assert "import pandas as pd" in source_text
    assert "import matplotlib.pyplot as plt" in source_text
    assert "read-only" in source_text.lower()
    assert "query_only" in source_text
    assert {
        "plot-run-outcomes",
        "plot-coverage-by-category",
        "plot-chunk-source-coverage",
        "plot-text-and-token-distributions",
        "plot-failure-reasons",
        "plot-failure-hotspots",
    }.issubset(cell_ids)
    assert all(not cell.get("outputs") for cell in notebook.cells)
    assert all(
        cell.get("execution_count") is None
        for cell in notebook.cells
        if cell.cell_type == "code"
    )


def _write_pdf(path: Path, text: str) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def _write_zero_page_pdf(path: Path) -> None:
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Count 0 /Kids [] >>\nendobj\n",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for pdf_object in objects:
        offsets.append(len(payload))
        payload.extend(pdf_object)
    xref_offset = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(payload)


def _write_blank_pdf(path: Path) -> None:
    import fitz

    document = fitz.open()
    document.new_page()
    document.save(path)
    document.close()


def _write_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("Document paragraph about cosine similarity.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "term"
    table.rows[0].cells[1].text = "weight"
    document.save(path)


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Retrieval Slide"
    slide.placeholders[1].text = "Slides can explain inverted indexes."
    presentation.save(path)


def _write_pptx_with_notes(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Retrieval Slide"
    slide.placeholders[1].text = "Slides can explain inverted indexes."
    slide.notes_slide.notes_text_frame.text = "Mention reciprocal rank fusion."
    presentation.save(path)


def _write_notebook(path: Path) -> None:
    notebook = nbformat.v4.new_notebook()
    notebook.cells = [
        nbformat.v4.new_markdown_cell("# Notebook Search Demo"),
        nbformat.v4.new_code_cell(
            "print('accuracy: 0.95')",
            outputs=[
                nbformat.v4.new_output(
                    output_type="stream",
                    name="stdout",
                    text="accuracy: 0.95\n",
                ),
                nbformat.v4.new_output(
                    output_type="display_data",
                    data={"image/png": "base64"},
                    metadata={},
                ),
            ],
            execution_count=1,
        ),
    ]
    nbformat.write(notebook, path)
